"""
Procesador de documentos para extraer texto de diferentes tipos de archivos.
"""
import os
import tempfile
from typing import Dict, Any, Tuple, TYPE_CHECKING
from core.utils import configurar_logging
from core.errors import ErrorProcesamientoDocumento

if TYPE_CHECKING:
    import PyPDF2
    import pptx
    import docx
    import pytesseract
    from PIL import Image
else:
    try:
        import PyPDF2
    except ImportError:
        PyPDF2 = None
    try:
        import pptx
    except ImportError:
        pptx = None
    try:
        import docx
    except ImportError:
        docx = None
    try:
        import pytesseract
        from PIL import Image
    except ImportError:
        pytesseract = None
        Image = None

class DocumentProcessor:
    """
    Procesador de documentos para extraer texto de diferentes tipos de archivos
    """
    def __init__(self):
        """
        Inicializa el procesador de documentos y verifica las dependencias disponibles.
        """
        self.logger = configurar_logging("document_processor")
        if PyPDF2 is None:
            self.logger.warning("PyPDF2 no está instalado. El procesamiento de PDF no estará disponible.")
        if pptx is None:
            self.logger.warning("python-pptx no está instalado. El procesamiento de PowerPoint no estará disponible.")
        if docx is None:
            self.logger.warning("python-docx no está instalado. El procesamiento de Word no estará disponible.")
        if pytesseract is None or Image is None:
            self.logger.warning("pytesseract o PIL no están instalados. El OCR no estará disponible.")

    def process_document(self, content: Any, content_type: str, filename: str) -> Tuple[str, Dict[str, Any]]:
        """
        Procesa un documento y extrae el texto según su tipo.
        
        Args:
            content: Contenido del archivo (bytes o str)
            content_type: Tipo MIME del archivo
            filename: Nombre del archivo
            
        Returns:
            Tupla con (texto extraído, metadatos)
            
        Raises:
            ErrorProcesamientoDocumento: Si ocurre un error al procesar el documento
        """
        metadata = {"filename": filename, "content_type": content_type}
        try:
            if content_type == "application/pdf" and PyPDF2:
                return self._process_pdf(content, metadata)
            elif content_type in ["application/vnd.openxmlformats-officedocument.presentationml.presentation", "application/vnd.ms-powerpoint"] and pptx:
                return self._process_pptx(content, metadata)
            elif content_type in ["application/vnd.openxmlformats-officedocument.wordprocessingml.document", "application/msword"] and docx:
                return self._process_docx(content, metadata)
            elif content_type.startswith("image/") and pytesseract and Image:
                return self._process_image(content, metadata)
            elif content_type.startswith("text/"):
                return str(content), metadata
            else:
                self.logger.warning(f"Tipo de archivo no soportado: {content_type}")
                return "", metadata
        except Exception as e:
            self.logger.error(f"Error al procesar documento: {e}")
            raise ErrorProcesamientoDocumento(str(e))

    def _process_pdf(self, content: bytes, metadata: Dict[str, Any]) -> Tuple[str, Dict[str, Any]]:
        """
        Extrae texto de un archivo PDF
        
        Args:
            content: Contenido binario del PDF
            metadata: Metadatos del documento
            
        Returns:
            Tupla con (texto extraído, metadatos actualizados)
            
        Raises:
            ErrorProcesamientoDocumento: Si ocurre un error al procesar el PDF
        """
        try:
            # Guardar el contenido en un archivo temporal
            with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as temp_file:
                temp_file.write(content)
                temp_path = temp_file.name

            # Extraer texto
            text = ""
            with open(temp_path, 'rb') as file:
                reader = PyPDF2.PdfReader(file)
                metadata['page_count'] = len(reader.pages)

                for page_num in range(len(reader.pages)):
                    page = reader.pages[page_num]
                    text += page.extract_text() + "\n\n"

            # Limpiar archivo temporal
            os.unlink(temp_path)

            return text, metadata

        except Exception as e:
            self.logger.error(f"Error al procesar PDF: {e}")
            raise ErrorProcesamientoDocumento(str(e))

    def _process_pptx(self, content: bytes, metadata: Dict[str, Any]) -> Tuple[str, Dict[str, Any]]:
        """
        Extrae texto de una presentación PowerPoint
        
        Args:
            content: Contenido binario de la presentación
            metadata: Metadatos del documento
            
        Returns:
            Tupla con (texto extraído, metadatos actualizados)
            
        Raises:
            ErrorProcesamientoDocumento: Si ocurre un error al procesar la presentación
        """
        try:
            # Guardar el contenido en un archivo temporal
            with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as temp_file:
                temp_file.write(content)
                temp_path = temp_file.name

            # Extraer texto
            text = ""
            presentation = pptx.Presentation(temp_path)
            metadata['slide_count'] = len(presentation.slides)

            for i, slide in enumerate(presentation.slides):
                text += f"Slide {i+1}:\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text += shape.text + "\n"
                text += "\n"

            # Limpiar archivo temporal
            os.unlink(temp_path)

            return text, metadata

        except Exception as e:
            self.logger.error(f"Error al procesar PowerPoint: {e}")
            raise ErrorProcesamientoDocumento(str(e))

    def _process_docx(self, content: bytes, metadata: Dict[str, Any]) -> Tuple[str, Dict[str, Any]]:
        """
        Extrae texto de un documento Word
        
        Args:
            content: Contenido binario del documento
            metadata: Metadatos del documento
            
        Returns:
            Tupla con (texto extraído, metadatos actualizados)
            
        Raises:
            ErrorProcesamientoDocumento: Si ocurre un error al procesar el documento
        """
        try:
            # Guardar el contenido en un archivo temporal
            with tempfile.NamedTemporaryFile(delete=False, suffix='.docx') as temp_file:
                temp_file.write(content)
                temp_path = temp_file.name

            # Extraer texto
            doc = docx.Document(temp_path)
            text = "\n".join([paragraph.text for paragraph in doc.paragraphs])

            # Limpiar archivo temporal
            os.unlink(temp_path)

            return text, metadata

        except Exception as e:
            self.logger.error(f"Error al procesar Word: {e}")
            raise ErrorProcesamientoDocumento(str(e))

    def _process_image(self, content: bytes, metadata: Dict[str, Any]) -> Tuple[str, Dict[str, Any]]:
        """
        Extrae texto de una imagen usando OCR
        
        Args:
            content: Contenido binario de la imagen
            metadata: Metadatos del documento
            
        Returns:
            Tupla con (texto extraído, metadatos actualizados)
            
        Raises:
            ErrorProcesamientoDocumento: Si ocurre un error al procesar la imagen
        """
        try:
            # Guardar el contenido en un archivo temporal
            with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as temp_file:
                temp_file.write(content)
                temp_path = temp_file.name

            # Extraer texto con OCR
            image = Image.open(temp_path)
            text = pytesseract.image_to_string(image)

            # Limpiar archivo temporal
            os.unlink(temp_path)

            return text, metadata

        except Exception as e:
            self.logger.error(f"Error al procesar imagen con OCR: {e}")
            raise ErrorProcesamientoDocumento(str(e))